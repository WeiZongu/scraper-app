"""
検索・抽出結果を、抽出キーワードを列とする表形式のPowerPointスライドとして
出力する。構成: 表紙スライド → 表スライド（1ページ5行、5行を超える分は
複数スライドに分割。ただし同じ出典ページ由来の行はなるべく同じページに
まとめる）。各行は抽出キーワードすべてを含む(AND条件)1件のレコードで、
列の値は該当キーワードを含む文・段落そのもの。各スライドの背景には、
そのページに含まれる行の画像から作った「フルカラースケルトン」
（輪郭線部分だけ元画像の色を残した透過画像）を敷く。
"""
from __future__ import annotations

import urllib.request
from datetime import datetime
from io import BytesIO
from pathlib import Path

from PIL import Image, ImageFilter, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from .config import SearchConfig

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

ACCENT_COLOR = RGBColor(0x30, 0x54, 0x96)
HEADER_TEXT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
ROW_BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
ROW_BG_ALT_COLOR = RGBColor(0xF0, 0xF3, 0xFA)
TEXT_COLOR = RGBColor(0x2B, 0x2B, 0x2B)
MUTED_COLOR = RGBColor(0x6B, 0x6B, 0x6B)
LINK_COLOR = RGBColor(0x1A, 0x4B, 0x9E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PLAIN_BG_COLOR = RGBColor(0xFA, 0xFB, 0xFD)  # 画像が無いページの背景

ROWS_PER_SLIDE = 5  # 1ページあたりの表示行数（ヘッダー行は別）
DOWNLOAD_TIMEOUT_SEC = 8
MAX_IMAGE_DIMENSION = 1280  # スケルトン加工前に縮小する上限(px)。メモリ節約のため


def _set_transparency(fill, alpha_pct: int) -> None:
    """python-pptxには塗りつぶし透明度の公開APIが無いため、XMLを直接操作する。
    alpha_pct: 0(完全に透明)〜100(完全に不透明)
    """
    color_elm = fill.fore_color._xFill.find(qn("a:srgbClr"))
    if color_elm is None:
        color_elm = fill.fore_color._xFill.find(qn("a:schemeClr"))
    if color_elm is None:
        return
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(int(alpha_pct * 1000)))
    color_elm.append(alpha)


def _download_image_bytes(url: str) -> BytesIO | None:
    if not url:
        return None
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=DOWNLOAD_TIMEOUT_SEC) as resp:
            data = resp.read()
        return BytesIO(data)
    except Exception:
        return None


def _load_image(url: str) -> Image.Image | None:
    """画像URLをダウンロードし、Pillowで扱える画像として読み込む。
    失敗した場合は None を返す。スケルトン加工の中間処理でのメモリ使用量を
    抑えるため、大きい画像はあらかじめ縮小する。
    """
    stream = _download_image_bytes(url)
    if stream is None:
        return None
    try:
        img = Image.open(stream)
        img.load()
        img = img.convert("RGB")
        img.thumbnail((MAX_IMAGE_DIMENSION, MAX_IMAGE_DIMENSION), Image.Resampling.LANCZOS)
        return img
    except Exception:
        return None


def _average_color(image: Image.Image) -> tuple[int, int, int]:
    small = image.resize((1, 1), Image.Resampling.LANCZOS)
    return small.getpixel((0, 0))


def _scale_color(rgb: tuple[int, int, int], factor: float) -> tuple[int, int, int]:
    return tuple(min(255, max(0, int(c * factor))) for c in rgb)


def _make_skeleton_overlay(image: Image.Image) -> Image.Image:
    """写真から輪郭線だけを抽出し、その部分だけ元画像のフルカラーを残した
    透過PNG(RGBA)を作る（「関連画像のスケルトン」の背景演出）。
    グレースケール化した画像に輪郭強調フィルタ(CONTOUR)をかけて輪郭線を
    抽出し、線の濃さをそのまま透明度に変換する。
    """
    gray = ImageOps.grayscale(image)
    blur_radius = max(1, min(gray.size) // 150)
    if blur_radius > 1:
        gray = gray.filter(ImageFilter.GaussianBlur(blur_radius))
    contour = gray.filter(ImageFilter.CONTOUR)

    alpha = ImageOps.invert(contour)
    alpha = ImageOps.autocontrast(alpha)
    alpha = alpha.point(lambda v: min(255, int(v * 1.6)))

    overlay = image.convert("RGBA")
    overlay.putalpha(alpha)
    return overlay


def _add_cover_picture(slide, pil_image: Image.Image) -> bool:
    """画像（PIL Image、透過PNGも可）をスライド全体を覆うように拡大配置する。
    失敗したら False。
    """
    try:
        buf = BytesIO()
        pil_image.save(buf, format="PNG")
        buf.seek(0)
        pic = slide.shapes.add_picture(buf, 0, 0)
        scale = max(SLIDE_WIDTH / pic.width, SLIDE_HEIGHT / pic.height)
        new_w = int(pic.width * scale)
        new_h = int(pic.height * scale)
        pic.width = new_w
        pic.height = new_h
        pic.left = int((SLIDE_WIDTH - new_w) / 2)
        pic.top = int((SLIDE_HEIGHT - new_h) / 2)
        return True
    except Exception:
        return False


def _add_solid_rect(slide, rgb: RGBColor, alpha_pct: int | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    if alpha_pct is not None:
        _set_transparency(shape.fill, alpha_pct)
    return shape


def _add_skeleton_background(slide, image_url: str) -> bool:
    """指定画像からフルカラースケルトンを作り、スライド全面の背景として敷く。
    画像の色合いから作った暗めの単色キャンバスの上に、輪郭線部分だけ元画像の
    色を残した透過画像を重ねる。成功したら True（呼び出し側で前景要素を
    不透明にする等、可読性を確保する必要がある）。
    """
    pil_image = _load_image(image_url)
    if pil_image is None:
        return False
    try:
        avg_color = _average_color(pil_image)
        canvas_color = RGBColor(*_scale_color(avg_color, 0.32))
        _add_solid_rect(slide, canvas_color)
        skeleton = _make_skeleton_overlay(pil_image)
        return _add_cover_picture(slide, skeleton)
    except Exception:
        return False


def _first_image_url(rows: list[dict]) -> str:
    for row in rows:
        if row.get("image_url"):
            return row["image_url"]
    return ""


def _add_title_slide(prs: Presentation, cfg: SearchConfig, rows: list[dict]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 白紙レイアウト
    has_bg = _add_skeleton_background(slide, _first_image_url(rows))
    if not has_bg:
        _add_solid_rect(slide, ACCENT_COLOR)
    else:
        # 背景の輪郭線が目立ちすぎて文字が読みにくくならないよう、軽く暗くする
        _add_solid_rect(slide, RGBColor(0x0A, 0x0E, 0x1A), 35)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.6), SLIDE_WIDTH - Inches(2), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = cfg.name or "検索結果"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    info_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), SLIDE_WIDTH - Inches(2), Inches(2))
    tf = info_box.text_frame
    tf.word_wrap = True
    lines = [
        f"検索キーワード: {cfg.search_keyword}",
        f"抽出キーワード（列）: {', '.join(cfg.extract_keywords)}",
        f"件数: {len(rows)}件　|　作成日時: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xD7, 0xDF, 0xF2)
        p.space_after = Pt(8)


def _set_cell_text(cell, text: str, *, size: int, color: RGBColor, bold: bool = False) -> None:
    tf = cell.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def _add_source_cell(cell, title: str, url: str, video_url: str, fetched_at: str) -> None:
    tf = cell.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR

    if url:
        p_url = tf.add_paragraph()
        run = p_url.add_run()
        run.text = url
        run.font.size = Pt(9)
        run.font.color.rgb = LINK_COLOR
        run.font.underline = True
        run.hyperlink.address = url

    if video_url:
        p_video = tf.add_paragraph()
        run = p_video.add_run()
        run.text = "[動画を見る]"
        run.font.size = Pt(9)
        run.font.color.rgb = LINK_COLOR
        run.font.underline = True
        run.hyperlink.address = video_url

    p_date = tf.add_paragraph()
    p_date.text = fetched_at
    p_date.font.size = Pt(8)
    p_date.font.color.rgb = MUTED_COLOR


def _add_table_slide(
    prs: Presentation,
    cfg: SearchConfig,
    rows_chunk: list[dict],
    page_num: int,
    total_pages: int,
) -> None:
    """抽出キーワードを列とする表を1枚のスライドに描画する（最大 ROWS_PER_SLIDE 行）。
    セル内のテキストは折り返し表示にする（複数行になった場合も枠内に収める）。
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    has_bg = _add_skeleton_background(slide, _first_image_url(rows_chunk))
    if has_bg:
        # 背景の輪郭線が目立ちすぎて表と重ならない上部の文字が読みにくくならないよう、軽く暗くする
        _add_solid_rect(slide, RGBColor(0x0A, 0x0E, 0x1A), 40)
        title_color = WHITE
        page_color = RGBColor(0xD7, 0xDF, 0xF2)
    else:
        _add_solid_rect(slide, PLAIN_BG_COLOR)
        title_color = ACCENT_COLOR
        page_color = MUTED_COLOR

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.6))
    tf = title_box.text_frame
    tf.text = cfg.name or "検索結果"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = title_color

    page_box = slide.shapes.add_textbox(SLIDE_WIDTH - Inches(1.8), Inches(0.3), Inches(1.3), Inches(0.4))
    page_box.text_frame.text = f"{page_num}/{total_pages}"
    page_box.text_frame.paragraphs[0].font.size = Pt(12)
    page_box.text_frame.paragraphs[0].font.color.rgb = page_color
    page_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    headers = list(cfg.extract_keywords) + ["出典"]
    n_cols = len(headers)
    n_rows_data = len(rows_chunk)
    n_rows = n_rows_data + 1  # ヘッダー行を含む

    table_left = Inches(0.5)
    table_top = Inches(1.1)
    table_width = SLIDE_WIDTH - Inches(1.0)
    table_height = SLIDE_HEIGHT - Inches(1.4)

    graphic_frame = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
    table = graphic_frame.table
    table.first_row = False
    table.horz_banding = False

    # 列幅: 「出典」列は少し広めに、残りをキーワード列で均等割りする
    source_col_width = Inches(2.6)
    keyword_col_width = int((table_width - source_col_width) / max(1, n_cols - 1))
    for c in range(n_cols - 1):
        table.columns[c].width = keyword_col_width
    table.columns[n_cols - 1].width = table_width - keyword_col_width * (n_cols - 1)

    header_row_height = Inches(0.5)
    data_row_height = int((table_height - header_row_height) / max(1, n_rows_data))
    table.rows[0].height = header_row_height
    for r in range(1, n_rows):
        table.rows[r].height = data_row_height

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        _set_cell_text(cell, header, size=13, color=HEADER_TEXT_COLOR, bold=True)

    for r, row in enumerate(rows_chunk, start=1):
        bg = ROW_BG_COLOR if r % 2 == 1 else ROW_BG_ALT_COLOR
        for c, kw in enumerate(cfg.extract_keywords):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            _set_cell_text(cell, row["columns"].get(kw, ""), size=12, color=TEXT_COLOR)

        source_cell = table.cell(r, n_cols - 1)
        source_cell.fill.solid()
        source_cell.fill.fore_color.rgb = bg
        _add_source_cell(
            source_cell,
            title=row.get("title", ""),
            url=row.get("url", ""),
            video_url=row.get("video_url", ""),
            fetched_at=row.get("fetched_at", ""),
        )


def _group_into_pages(rows: list[dict], rows_per_page: int) -> list[list[dict]]:
    """同じ出典ページ（url）から抽出された行＝関連性の高いデータをなるべく
    同じスライドにまとめてから、1ページあたりの行数上限でページに割り付ける。
    グループがページの残り行数に収まらない場合は、無駄な分断を避けるため
    新しいページから始める（グループ自体がページ上限より大きい場合のみ、
    やむを得ず複数ページにまたがる）。
    """
    groups: dict[str, list[dict]] = {}
    order: list[str] = []
    for row in rows:
        key = row.get("url") or row.get("title") or ""
        if key not in groups:
            groups[key] = []
            order.append(key)
        groups[key].append(row)

    pages: list[list[dict]] = []
    current: list[dict] = []
    for key in order:
        group = groups[key]
        if len(group) <= rows_per_page and len(current) + len(group) > rows_per_page:
            pages.append(current)
            current = []
        for row in group:
            if len(current) >= rows_per_page:
                pages.append(current)
                current = []
            current.append(row)
    if current:
        pages.append(current)
    return pages


def export_report(cfg: SearchConfig, rows: list[dict], output_path: Path) -> Path:
    """
    rows: engine.scrape() が返す辞書のリスト
          （columns: {キーワード: 値, ...} / image_url / video_url / title / url /
          fetched_at の各キーを持つ）
    表紙 → 表スライド（1ページ ROWS_PER_SLIDE 行、同じ出典ページの行はなるべく
    同じスライドにまとめる）の順に構成する。
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _add_title_slide(prs, cfg, rows)

    pages = _group_into_pages(rows, ROWS_PER_SLIDE)
    total_pages = max(1, len(pages))
    for page_num, chunk in enumerate(pages, start=1):
        _add_table_slide(prs, cfg, chunk, page_num, total_pages)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
